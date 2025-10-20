# DocSwap - Flask Backend Server

from flask import Flask, request, jsonify, send_file, render_template, url_for
import os
import uuid
import time
import json
from werkzeug.utils import secure_filename
import shutil
from dotenv import load_dotenv
import jwt
import requests
from datetime import datetime, timezone
import re
from flask_cors import CORS
from flask_limiter import Limiter
from flask_limiter.util import get_remote_address
import logging
import threading
import schedule
import smtplib
from email.mime.text import MIMEText
from email.mime.multipart import MIMEMultipart

# Load environment variables
load_dotenv()

# =============================================================================
# UNIFIED PRODUCTION CONFIGURATION
# =============================================================================
# This app.py serves as the single unified production build for DocSwap
# It supports deployment to multiple servers:
# - mydocswap.com (primary domain)
# - 31.97.193.96 (backup/development server)
# 
# The app includes:
# - Enhanced conversion system with intelligent fallback
# - Comprehensive error handling and logging
# - Security headers and rate limiting
# - Multi-server CORS support
# =============================================================================

# Enhanced conversion system
try:
    from conversion import ConversionManager
    conversion_system_available = True
except ImportError as e:
    print(f"Enhanced conversion system not available: {e}")
    ConversionManager = None
    conversion_system_available = False

# Import conversion libraries
try:
    from supabase import create_client, Client
    from pdf2docx import Converter as PDFToDocx
    import fitz  # PyMuPDF for PDF to image
    import pytesseract  # OCR
    from PIL import Image
    import pandas as pd
    import pptx
    import docx2txt
    import img2pdf
    import docx2pdf
    import PyJWT
    PIL_AVAILABLE = True
except ImportError as e:
    print("Some conversion libraries are missing. Please install them using:")
    print("pip install pdf2docx PyMuPDF pytesseract pillow pandas python-pptx img2pdf docx2pdf schedule")
    print(f"Import error: {e}")
    PIL_AVAILABLE = False
    Image = None

# Try to import python-magic for file content validation
try:
    import magic
    MAGIC_AVAILABLE = True
except ImportError:
    MAGIC_AVAILABLE = False
    magic = None

app = Flask(__name__, static_folder='.', static_url_path='')

# Track app start time for uptime monitoring
app.start_time = time.time()

# Configure Flask app
app.config['SECRET_KEY'] = os.environ.get('SECRET_KEY', 'dev-key-change-in-production')
app.config['MAX_CONTENT_LENGTH'] = int(os.environ.get('MAX_FILE_SIZE', 104857600))  # 100MB default

# CORS configuration - Unified for all deployment targets
default_origins = 'https://mydocswap.com,https://www.mydocswap.com,http://31.97.193.96:8000,https://31.97.193.96:8000,http://localhost:8000,http://127.0.0.1:8000'
allowed_origins = os.environ.get('ALLOWED_ORIGINS', default_origins).split(',')
CORS(app, origins=allowed_origins)

# Security headers middleware
@app.after_request
def add_security_headers(response):
    """Add security headers to all responses"""
    # Prevent clickjacking
    response.headers['X-Frame-Options'] = 'DENY'
    
    # Prevent MIME type sniffing
    response.headers['X-Content-Type-Options'] = 'nosniff'
    
    # Enable XSS protection
    response.headers['X-XSS-Protection'] = '1; mode=block'
    
    # Strict Transport Security (HTTPS only)
    if request.is_secure:
        response.headers['Strict-Transport-Security'] = 'max-age=31536000; includeSubDomains'
    
    # Content Security Policy - Comprehensive configuration for all external resources
    csp = (
        "default-src 'self'; "
        "script-src 'self' 'unsafe-inline' 'unsafe-eval' https://cdn.jsdelivr.net https://unpkg.com; "
        "style-src 'self' 'unsafe-inline' https://fonts.googleapis.com; "
        "img-src 'self' data: blob:; "
        "font-src 'self' https://fonts.gstatic.com; "
        "connect-src 'self' *.supabase.co qzuwonueyvouvrhiwcob.supabase.co; "
        "media-src 'self'; "
        "object-src 'none'; "
        "child-src 'none'; "
        "worker-src 'none'; "
        "frame-ancestors 'self'; "
        "form-action 'self'; "
        "base-uri 'self'"
    )
    response.headers['Content-Security-Policy'] = csp
    
    # Referrer Policy
    response.headers['Referrer-Policy'] = 'strict-origin-when-cross-origin'
    
    # Permissions Policy
    response.headers['Permissions-Policy'] = (
        "geolocation=(), "
        "microphone=(), "
        "camera=(), "
        "payment=(), "
        "usb=(), "
        "magnetometer=(), "
        "gyroscope=(), "
        "speaker=()"
    )
    
    return response

# Rate limiting
limiter = Limiter(
    key_func=get_remote_address,
    default_limits=[f"{os.environ.get('RATE_LIMIT_PER_MINUTE', 60)} per minute"]
)
limiter.init_app(app)

# Rate limiting error handler
@app.errorhandler(429)
def ratelimit_handler(e):
    return jsonify({
        'error': 'Rate limit exceeded',
        'message': 'Too many requests. Please wait before trying again.',
        'retry_after': '60'
    }), 429

# Configure logging
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s %(levelname)s %(name)s %(message)s'
)
logger = logging.getLogger(__name__)

# Supabase authentication configuration
SUPABASE_URL = os.environ.get('SUPABASE_URL')
SUPABASE_SERVICE_KEY = os.environ.get('SUPABASE_SERVICE_KEY')
SUPABASE_ANON_KEY = os.environ.get('SUPABASE_ANON_KEY')
SUPABASE_JWT_SECRET = os.environ.get('SUPABASE_JWT_SECRET')

if not all([SUPABASE_URL, SUPABASE_SERVICE_KEY, SUPABASE_ANON_KEY, SUPABASE_JWT_SECRET]):
    logger.error("Missing required Supabase environment variables")
    raise ValueError("Missing required Supabase environment variables")

# Configuration
UPLOAD_FOLDER = os.path.join(os.path.dirname(os.path.abspath(__file__)), 'uploads')
OUTPUT_FOLDER = os.path.join(os.path.dirname(os.path.abspath(__file__)), 'output')
ALLOWED_EXTENSIONS = {'pdf', 'docx', 'doc', 'xlsx', 'xls', 'pptx', 'ppt', 'jpg', 'jpeg', 'png', 'webp', 'bmp', 'tiff', 'gif', 'svg', 'txt', 'html', 'csv', 'json'}
MAX_FILE_SIZE = int(os.environ.get('MAX_FILE_SIZE', 104857600))  # 100MB default
FILE_EXPIRY = int(os.environ.get('FILE_EXPIRY', 86400))  # 24 hours default

# Security patterns for input validation
FILENAME_PATTERN = re.compile(r'^[a-zA-Z0-9._-]+$')
EMAIL_PATTERN = re.compile(r'^[a-zA-Z0-9._%+-]+@[a-zA-Z0-9.-]+\.[a-zA-Z]{2,}$')

# Helper function to get remote address
def get_remote_address():
    """Get the remote IP address from request headers or direct connection"""
    # Check for forwarded headers (common in production with reverse proxies)
    if 'X-Forwarded-For' in request.headers:
        return request.headers['X-Forwarded-For'].split(',')[0].strip()
    elif 'X-Real-IP' in request.headers:
        return request.headers['X-Real-IP']
    else:
        return request.remote_addr or 'unknown'

# Create folders if they don't exist
os.makedirs(UPLOAD_FOLDER, exist_ok=True)
os.makedirs(OUTPUT_FOLDER, exist_ok=True)

# Session storage directory
SESSIONS_FOLDER = os.path.join(os.path.dirname(os.path.abspath(__file__)), 'sessions')
os.makedirs(SESSIONS_FOLDER, exist_ok=True)

# Persistent session storage class
class PersistentSessionStorage:
    def __init__(self, sessions_folder):
        self.sessions_folder = sessions_folder
        
    def _get_session_file(self, session_id):
        return os.path.join(self.sessions_folder, f"{session_id}.json")
    
    def get(self, session_id, default=None):
        session_file = self._get_session_file(session_id)
        if os.path.exists(session_file):
            try:
                with open(session_file, 'r') as f:
                    return json.load(f)
            except (json.JSONDecodeError, IOError) as e:
                logger.warning(f"Failed to load session {session_id}: {e}")
                return default
        return default
    
    def __getitem__(self, session_id):
        session = self.get(session_id)
        if session is None:
            raise KeyError(session_id)
        return session
    
    def __setitem__(self, session_id, session_data):
        session_file = self._get_session_file(session_id)
        try:
            with open(session_file, 'w') as f:
                json.dump(session_data, f, indent=2)
        except IOError as e:
            logger.error(f"Failed to save session {session_id}: {e}")
    
    def __contains__(self, session_id):
        return os.path.exists(self._get_session_file(session_id))
    
    def __delitem__(self, session_id):
        session_file = self._get_session_file(session_id)
        if os.path.exists(session_file):
            try:
                os.remove(session_file)
            except IOError as e:
                logger.warning(f"Failed to delete session {session_id}: {e}")
    
    def keys(self):
        """Return all session IDs"""
        session_files = []
        for filename in os.listdir(self.sessions_folder):
            if filename.endswith('.json'):
                session_files.append(filename[:-5])  # Remove .json extension
        return session_files
    
    def items(self):
        """Return all session items"""
        for session_id in self.keys():
            session_data = self.get(session_id)
            if session_data is not None:
                yield session_id, session_data

# Initialize persistent session storage
sessions = PersistentSessionStorage(SESSIONS_FOLDER)

# Initialize enhanced conversion manager
if conversion_system_available and ConversionManager:
    try:
        conversion_manager = ConversionManager()
        logger.info("Enhanced conversion system initialized successfully")
    except Exception as e:
        logger.error(f"Failed to initialize conversion system: {str(e)}")
        conversion_manager = None
else:
    conversion_manager = None
    logger.info("Enhanced conversion system not available, using fallback methods")

# Helper functions
def sanitize_filename(filename):
    """Sanitize filename to prevent path traversal and other attacks"""
    if not filename:
        return None
    
    # Remove any path components
    filename = os.path.basename(filename)
    
    # Check if filename matches safe pattern
    if not FILENAME_PATTERN.match(filename.replace(' ', '_')):
        # If not safe, create a safe version
        safe_chars = re.sub(r'[^a-zA-Z0-9._-]', '_', filename)
        filename = safe_chars
    
    # Limit filename length
    if len(filename) > 255:
        name, ext = os.path.splitext(filename)
        filename = name[:250] + ext
    
    return filename

def allowed_file(filename):
    """Check if file extension is allowed"""
    if not filename or '.' not in filename:
        return False
    
    extension = filename.rsplit('.', 1)[1].lower()
    return extension in ALLOWED_EXTENSIONS

def validate_file_content(file_path, expected_extension):
    """Validate file content matches expected type"""
    if not MAGIC_AVAILABLE:
        # python-magic is not available, skip content validation
        return {'valid': True}
        
    try:
        mime = magic.Magic(mime=True)
        file_mime = mime.from_file(file_path)
        
        # Define expected MIME types for extensions
        mime_map = {
            'pdf': 'application/pdf',
            'docx': 'application/vnd.openxmlformats-officedocument.wordprocessingml.document',
            'jpg': 'image/jpeg',
            'jpeg': 'image/jpeg',
            'png': 'image/png',
            'txt': 'text/plain'
        }
        
        expected_mime = mime_map.get(expected_extension.lower())
        if expected_mime and not file_mime.startswith(expected_mime.split('/')[0]):
            return {'valid': False, 'error': f'File content does not match expected type {expected_extension}'}
            
    except Exception as e:
        logger.warning(f"File content validation failed: {str(e)}")
        return {'valid': False, 'error': f'File validation error: {str(e)}'}
    
    return {'valid': True}

def get_file_extension(filename):
    """Get file extension safely"""
    if not filename or '.' not in filename:
        return ''
    return filename.rsplit('.', 1)[1].lower()

def generate_unique_filename(original_filename, output_format):
    """Generate unique filename with proper sanitization"""
    # Sanitize the original filename
    safe_filename = sanitize_filename(original_filename)
    if not safe_filename:
        safe_filename = "file"
    
    # Get the filename without extension
    name = safe_filename.rsplit('.', 1)[0] if '.' in safe_filename else safe_filename
    
    # Generate a unique ID
    unique_id = str(uuid.uuid4())[:8]
    
    # Return the new filename
    return f"{name}-{unique_id}.{output_format}"

# Conversion functions
def convert_pdf_to_docx(input_path, output_path):
    converter = PDFToDocx(input_path)
    converter.convert(output_path)
    converter.close()
    return output_path

def convert_pdf_to_image(input_path, output_path, format='jpg'):
    try:
        pdf_document = fitz.open(input_path)
        # Just convert the first page for preview
        page = pdf_document.load_page(0)
        
        # Ensure the output path has the correct extension
        if not output_path.lower().endswith(f'.{format.lower()}'):
            output_path = f"{os.path.splitext(output_path)[0]}.{format.lower()}"
        
        # Create a pixmap with lower zoom factor for faster processing
        # Reduced from 2.0 to 1.5 for better performance while maintaining decent quality
        zoom = 1.2  # Lower zoom factor for faster processing
        matrix = fitz.Matrix(zoom, zoom)
        pix = page.get_pixmap(matrix=matrix, alpha=False)
        
        # Direct save to output format when possible
        if format.lower() in ['png']:
            # Save directly to PNG for better performance
            pix.save(output_path)
            return output_path
        
        # For other formats, use PIL with optimized settings
        from PIL import Image
        import io
        
        # Convert pixmap to PIL Image using memory buffer for efficiency
        img_data = pix.tobytes("png")
        img = Image.open(io.BytesIO(img_data))
        
        # Optimize based on format
        if format.lower() in ['jpg', 'jpeg']:
            img = img.convert('RGB')
            # Lower quality for faster processing (85 instead of 95)
            img.save(output_path, format='JPEG', quality=85, optimize=True)
        elif format.lower() == 'png':
            img.save(output_path, format='PNG', optimize=True, compress_level=6)  # Lower compression level
        
        return output_path
        
        # Verify the file was created successfully
        if not os.path.exists(output_path) or os.path.getsize(output_path) == 0:
            raise Exception(f"Failed to create image file at {output_path}")
            
        print(f"Successfully converted PDF to {format} image: {output_path}")
        return output_path
    except Exception as e:
        print(f"Error in convert_pdf_to_image: {str(e)}")
        raise

def convert_pdf_to_text(input_path, output_path, use_ocr=False):
    if use_ocr:
        # Use OCR to extract text
        pdf_document = fitz.open(input_path)
        text = ""
        for page_num in range(len(pdf_document)):
            page = pdf_document.load_page(page_num)
            pix = page.get_pixmap()
            img_path = f"{output_path}.temp.png"
            pix.save(img_path)
            text += pytesseract.image_to_string(Image.open(img_path)) + "\n\n"
            os.remove(img_path)  # Clean up temp image
    else:
        # Extract text directly from PDF
        pdf_document = fitz.open(input_path)
        text = ""
        for page_num in range(len(pdf_document)):
            page = pdf_document.load_page(page_num)
            text += page.get_text() + "\n\n"
    
    # Write text to output file
    with open(output_path, 'w', encoding='utf-8') as f:
        f.write(text)
    
    return output_path

def convert_image_to_pdf(input_path, output_path):
    """Convert image to PDF with enhanced error handling"""
    try:
        # Validate input image
        with Image.open(input_path) as img:
            # Convert to RGB if necessary (for JPEG compatibility)
            if img.mode in ('RGBA', 'LA', 'P'):
                img = img.convert('RGB')
            
            # Save as temporary JPEG if needed for img2pdf compatibility
            temp_path = input_path
            if img.format not in ['JPEG', 'PNG']:
                temp_path = f"{input_path}.temp.jpg"
                img.save(temp_path, 'JPEG', quality=95)
        
        # Convert to PDF
        with open(output_path, "wb") as f:
            f.write(img2pdf.convert(temp_path))
        
        # Clean up temporary file if created
        if temp_path != input_path and os.path.exists(temp_path):
            os.remove(temp_path)
            
        return output_path
    except Exception as e:
        print(f"Error in convert_image_to_pdf: {str(e)}")
        raise

def convert_image_to_image(input_path, output_path, output_format):
    """Convert between different image formats"""
    try:
        with Image.open(input_path) as img:
            # Handle transparency for formats that don't support it
            if output_format.lower() in ['jpg', 'jpeg'] and img.mode in ('RGBA', 'LA'):
                # Create white background for JPEG
                background = Image.new('RGB', img.size, (255, 255, 255))
                if img.mode == 'RGBA':
                    background.paste(img, mask=img.split()[-1])  # Use alpha channel as mask
                else:
                    background.paste(img)
                img = background
            elif output_format.lower() == 'png' and img.mode not in ('RGBA', 'LA', 'P'):
                img = img.convert('RGBA')
            
            # Save with appropriate settings
            save_kwargs = {}
            if output_format.lower() in ['jpg', 'jpeg']:
                save_kwargs = {'quality': 95, 'optimize': True}
                img = img.convert('RGB')
            elif output_format.lower() == 'png':
                save_kwargs = {'optimize': True}
            elif output_format.lower() == 'webp':
                save_kwargs = {'quality': 95, 'method': 6}
            
            img.save(output_path, format=output_format.upper(), **save_kwargs)
        return output_path
    except Exception as e:
        print(f"Error in convert_image_to_image: {str(e)}")
        raise

def convert_docx_to_pdf(input_path, output_path):
    """Convert DOCX to PDF with enhanced error handling"""
    try:
        docx2pdf.convert(input_path, output_path)
        return output_path
    except Exception as e:
        print(f"Error in convert_docx_to_pdf: {str(e)}")
        raise

def convert_docx_to_text(input_path, output_path):
    """Convert DOCX to plain text"""
    try:
        from docx import Document
        doc = Document(input_path)
        text = []
        for paragraph in doc.paragraphs:
            text.append(paragraph.text)
        
        with open(output_path, 'w', encoding='utf-8') as f:
            f.write('\n'.join(text))
        return output_path
    except Exception as e:
        print(f"Error in convert_docx_to_text: {str(e)}")
        raise

def convert_text_to_pdf(input_path, output_path):
    """Convert text file to PDF"""
    try:
        from reportlab.pdfgen import canvas
        from reportlab.lib.pagesizes import letter
        
        c = canvas.Canvas(output_path, pagesize=letter)
        width, height = letter
        
        with open(input_path, 'r', encoding='utf-8') as f:
            lines = f.readlines()
        
        y = height - 50  # Start from top
        line_height = 14
        
        for line in lines:
            if y < 50:  # Start new page if near bottom
                c.showPage()
                y = height - 50
            
            # Handle long lines by wrapping
            line = line.rstrip('\n\r')
            if len(line) > 80:
                words = line.split(' ')
                current_line = ''
                for word in words:
                    if len(current_line + word) < 80:
                        current_line += word + ' '
                    else:
                        c.drawString(50, y, current_line.strip())
                        y -= line_height
                        current_line = word + ' '
                        if y < 50:
                            c.showPage()
                            y = height - 50
                if current_line.strip():
                    c.drawString(50, y, current_line.strip())
                    y -= line_height
            else:
                c.drawString(50, y, line)
                y -= line_height
        
        c.save()
        return output_path
    except Exception as e:
        print(f"Error in convert_text_to_pdf: {str(e)}")
        raise

def convert_excel_to_pdf(input_path, output_path):
    """Convert Excel file to PDF"""
    try:
        import pandas as pd
        from reportlab.pdfgen import canvas
        from reportlab.lib.pagesizes import letter, landscape
        from reportlab.lib import colors
        from reportlab.platypus import SimpleDocTemplate, Table, TableStyle
        
        # Read Excel file
        df = pd.read_excel(input_path, sheet_name=None)  # Read all sheets
        
        # Create PDF
        doc = SimpleDocTemplate(output_path, pagesize=landscape(letter))
        elements = []
        
        for sheet_name, sheet_df in df.items():
            # Add sheet title
            from reportlab.platypus import Paragraph, Spacer
            from reportlab.lib.styles import getSampleStyleSheet
            styles = getSampleStyleSheet()
            
            title = Paragraph(f"<b>Sheet: {sheet_name}</b>", styles['Heading1'])
            elements.append(title)
            elements.append(Spacer(1, 12))
            
            # Convert DataFrame to table data
            data = [sheet_df.columns.tolist()]  # Header
            data.extend(sheet_df.values.tolist())  # Data rows
            
            # Limit data size for PDF
            if len(data) > 100:  # Limit rows
                data = data[:100]
            
            # Limit column width
            max_cols = 10
            if len(data[0]) > max_cols:
                data = [row[:max_cols] for row in data]
            
            # Create table
            table = Table(data)
            table.setStyle(TableStyle([
                ('BACKGROUND', (0, 0), (-1, 0), colors.grey),
                ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
                ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
                ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
                ('FONTSIZE', (0, 0), (-1, 0), 8),
                ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
                ('BACKGROUND', (0, 1), (-1, -1), colors.beige),
                ('FONTSIZE', (0, 1), (-1, -1), 6),
                ('GRID', (0, 0), (-1, -1), 1, colors.black)
            ]))
            
            elements.append(table)
            elements.append(Spacer(1, 20))
        
        doc.build(elements)
        return output_path
    except Exception as e:
        print(f"Error in convert_excel_to_pdf: {str(e)}")
        raise

def convert_excel_to_csv(input_path, output_path):
    """Convert Excel file to CSV (first sheet only)"""
    try:
        import pandas as pd
        
        # Read first sheet of Excel file
        df = pd.read_excel(input_path, sheet_name=0)
        
        # Save as CSV
        df.to_csv(output_path, index=False)
        return output_path
    except Exception as e:
        print(f"Error in convert_excel_to_csv: {str(e)}")
        raise

def convert_excel_to_html(input_path, output_path):
    """Convert Excel file to HTML"""
    try:
        import pandas as pd
        
        # Read all sheets
        df_dict = pd.read_excel(input_path, sheet_name=None)
        
        html_content = """
        <!DOCTYPE html>
        <html>
        <head>
            <title>Excel Data</title>
            <style>
                body { font-family: Arial, sans-serif; margin: 20px; }
                table { border-collapse: collapse; width: 100%; margin-bottom: 20px; }
                th, td { border: 1px solid #ddd; padding: 8px; text-align: left; }
                th { background-color: #f2f2f2; font-weight: bold; }
                h2 { color: #333; border-bottom: 2px solid #333; }
            </style>
        </head>
        <body>
        """
        
        for sheet_name, df in df_dict.items():
            html_content += f"<h2>Sheet: {sheet_name}</h2>\n"
            html_content += df.to_html(index=False, escape=False, classes='excel-table')
            html_content += "\n"
        
        html_content += """
        </body>
        </html>
        """
        
        with open(output_path, 'w', encoding='utf-8') as f:
            f.write(html_content)
        
        return output_path
    except Exception as e:
        print(f"Error in convert_excel_to_html: {str(e)}")
        raise

def convert_powerpoint_to_pdf(input_path, output_path):
    """Convert PowerPoint file to PDF"""
    try:
        from pptx import Presentation
        from reportlab.pdfgen import canvas
        from reportlab.lib.pagesizes import letter, landscape
        from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Image as RLImage
        from reportlab.lib.styles import getSampleStyleSheet
        from reportlab.lib.units import inch
        import tempfile
        import os
        
        # Load presentation
        prs = Presentation(input_path)
        
        # Create PDF
        doc = SimpleDocTemplate(output_path, pagesize=landscape(letter))
        elements = []
        styles = getSampleStyleSheet()
        
        for i, slide in enumerate(prs.slides):
            # Add slide title
            title = Paragraph(f"<b>Slide {i + 1}</b>", styles['Heading1'])
            elements.append(title)
            elements.append(Spacer(1, 12))
            
            # Extract text from slide
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text.strip())
            
            if slide_text:
                for text in slide_text:
                    para = Paragraph(text, styles['Normal'])
                    elements.append(para)
                    elements.append(Spacer(1, 6))
            else:
                para = Paragraph("(No text content)", styles['Italic'])
                elements.append(para)
            
            elements.append(Spacer(1, 20))
        
        doc.build(elements)
        return output_path
    except Exception as e:
        print(f"Error in convert_powerpoint_to_pdf: {str(e)}")
        raise

def convert_powerpoint_to_images(input_path, output_path, format='png'):
    """Convert PowerPoint slides to images"""
    try:
        from pptx import Presentation
        from PIL import Image, ImageDraw, ImageFont
        import os
        
        # Load presentation
        prs = Presentation(input_path)
        
        # Create a simple image for each slide with text
        images = []
        for i, slide in enumerate(prs.slides):
            # Create a blank image
            img = Image.new('RGB', (1024, 768), color='white')
            draw = ImageDraw.Draw(img)
            
            # Try to use a default font
            try:
                font = ImageFont.truetype("/System/Library/Fonts/Arial.ttf", 24)
                title_font = ImageFont.truetype("/System/Library/Fonts/Arial.ttf", 32)
            except:
                font = ImageFont.load_default()
                title_font = ImageFont.load_default()
            
            # Add slide number
            draw.text((20, 20), f"Slide {i + 1}", fill='black', font=title_font)
            
            # Extract and draw text from slide
            y_position = 80
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text = shape.text.strip()
                    # Wrap text if too long
                    words = text.split()
                    lines = []
                    current_line = []
                    for word in words:
                        current_line.append(word)
                        if len(' '.join(current_line)) > 80:  # Approximate line length
                            if len(current_line) > 1:
                                lines.append(' '.join(current_line[:-1]))
                                current_line = [word]
                            else:
                                lines.append(word)
                                current_line = []
                    if current_line:
                        lines.append(' '.join(current_line))
                    
                    for line in lines:
                        if y_position < 700:  # Don't go off the image
                            draw.text((20, y_position), line, fill='black', font=font)
                            y_position += 30
            
            images.append(img)
        
        # Save images
        if len(images) == 1:
            # Single slide - save directly
            images[0].save(output_path, format.upper())
        else:
            # Multiple slides - save as separate files or combine
            base_name = os.path.splitext(output_path)[0]
            for i, img in enumerate(images):
                slide_path = f"{base_name}_slide_{i+1}.{format.lower()}"
                img.save(slide_path, format.upper())
            
            # Also save the first slide as the main output
            images[0].save(output_path, format.upper())
        
        return output_path
    except Exception as e:
        print(f"Error in convert_powerpoint_to_images: {str(e)}")
        raise

def convert_csv_to_excel(input_path, output_path):
    """Convert CSV file to Excel"""
    try:
        import pandas as pd
        
        # Read CSV file
        df = pd.read_csv(input_path)
        
        # Save as Excel
        df.to_excel(output_path, index=False)
        return output_path
    except Exception as e:
        print(f"Error in convert_csv_to_excel: {str(e)}")
        raise

def convert_pdf_to_html(input_path, output_path):
    """Convert PDF to HTML"""
    try:
        pdf_document = fitz.open(input_path)
        html_content = ['<html><head><title>Converted PDF</title></head><body>']
        
        for page_num in range(len(pdf_document)):
            page = pdf_document.load_page(page_num)
            text = page.get_text()
            # Simple HTML formatting
            paragraphs = text.split('\n\n')
            for para in paragraphs:
                if para.strip():
                    html_content.append(f'<p>{para.strip()}</p>')
        
        html_content.append('</body></html>')
        
        with open(output_path, 'w', encoding='utf-8') as f:
            f.write('\n'.join(html_content))
        
        return output_path
    except Exception as e:
        print(f"Error in convert_pdf_to_html: {str(e)}")
        raise

def convert_text_to_html(input_path, output_path):
    """Convert text file to HTML"""
    try:
        with open(input_path, 'r', encoding='utf-8') as f:
            text = f.read()
        
        html_content = ['<html><head><title>Converted Text</title></head><body>']
        
        # Convert line breaks to paragraphs
        paragraphs = text.split('\n\n')
        for para in paragraphs:
            if para.strip():
                # Replace single line breaks with <br> tags
                formatted_para = para.replace('\n', '<br>')
                html_content.append(f'<p>{formatted_para}</p>')
        
        html_content.append('</body></html>')
        
        with open(output_path, 'w', encoding='utf-8') as f:
            f.write('\n'.join(html_content))
        
        return output_path
    except Exception as e:
        print(f"Error in convert_text_to_html: {str(e)}")
        raise

def convert_html_to_pdf(input_path, output_path):
    """Convert HTML to PDF"""
    try:
        import pdfkit
        pdfkit.from_file(input_path, output_path)
        return output_path
    except ImportError:
        # Fallback using reportlab for basic HTML
        try:
            from reportlab.pdfgen import canvas
            from reportlab.lib.pagesizes import letter
            import html
            
            with open(input_path, 'r', encoding='utf-8') as f:
                html_content = f.read()
            
            # Strip HTML tags for basic conversion
            import re
            text = re.sub('<[^<]+?>', '', html_content)
            text = html.unescape(text)
            
            c = canvas.Canvas(output_path, pagesize=letter)
            width, height = letter
            
            lines = text.split('\n')
            y = height - 50
            line_height = 14
            
            for line in lines:
                if y < 50:
                    c.showPage()
                    y = height - 50
                c.drawString(50, y, line[:80])  # Limit line length
                y -= line_height
            
            c.save()
            return output_path
        except Exception as e:
            print(f"Error in convert_html_to_pdf fallback: {str(e)}")
            raise
    except Exception as e:
        print(f"Error in convert_html_to_pdf: {str(e)}")
        raise

def convert_html_to_text(input_path, output_path):
    """Convert HTML to plain text"""
    try:
        import html
        import re
        
        with open(input_path, 'r', encoding='utf-8') as f:
            html_content = f.read()
        
        # Strip HTML tags
        text = re.sub('<[^<]+?>', '', html_content)
        text = html.unescape(text)
        
        # Clean up extra whitespace
        text = re.sub(r'\n\s*\n', '\n\n', text)
        text = text.strip()
        
        with open(output_path, 'w', encoding='utf-8') as f:
            f.write(text)
        
        return output_path
    except Exception as e:
        print(f"Error in convert_html_to_text: {str(e)}")
        raise

# Verify Supabase JWT token with proper signature validation
def verify_supabase_token(request):
    """Verify Supabase JWT token with proper signature validation"""
    auth_header = request.headers.get('Authorization')
    if not auth_header or not auth_header.startswith('Bearer '):
        logger.warning("Missing or invalid Authorization header")
        return None
    
    token = auth_header.split(' ')[1]
    
    try:
        # Decode and verify the JWT token
        decoded_token = jwt.decode(
            token,
            SUPABASE_JWT_SECRET,
            algorithms=['HS256'],
            audience='authenticated',
            issuer='supabase'
        )
        
        # Check token expiration
        exp = decoded_token.get('exp')
        if exp and datetime.fromtimestamp(exp, tz=timezone.utc) < datetime.now(tz=timezone.utc):
            logger.warning("Token has expired")
            return None
        
        # Validate required claims
        required_claims = ['sub', 'email', 'aud']
        for claim in required_claims:
            if claim not in decoded_token:
                logger.warning(f"Missing required claim: {claim}")
                return None
        
        # Validate email format
        email = decoded_token.get('email')
        if email and not EMAIL_PATTERN.match(email):
            logger.warning("Invalid email format in token")
            return None
        
        logger.info(f"Token verified successfully for user: {decoded_token.get('email')}")
        return decoded_token
        
    except jwt.ExpiredSignatureError:
        logger.warning("Token has expired")
        return None
    except jwt.InvalidTokenError as e:
        logger.warning(f"Invalid token: {str(e)}")
        return None
    except Exception as e:
        logger.error(f"Token verification error: {str(e)}")
        return None

# Example of a protected route that requires authentication
@app.route('/api/user/profile', methods=['GET'])
def user_profile():
    user_data = verify_supabase_token(request)
    if not user_data:
        return jsonify({'error': 'Unauthorized'}), 401
    
    # Return user profile data
    # In a real application, you might fetch additional user data from a database
    return jsonify({
        'user_id': user_data.get('sub'),
        'email': user_data.get('email'),
        'auth_provider': user_data.get('aud')
    })

# Cleanup expired files
def cleanup_expired_files():
    current_time = time.time()
    # Set cleanup time to 24 hours (86400 seconds)
    cleanup_time = 24 * 3600  # 24 hours in seconds
    
    cleaned_files = []
    
    # Check uploads folder
    if os.path.exists(UPLOAD_FOLDER):
        for filename in os.listdir(UPLOAD_FOLDER):
            file_path = os.path.join(UPLOAD_FOLDER, filename)
            # If file is older than 24 hours, delete it
            if os.path.isfile(file_path) and current_time - os.path.getmtime(file_path) > cleanup_time:
                try:
                    os.remove(file_path)
                    cleaned_files.append(f"uploads/{filename}")
                    logger.info(f"Auto cleanup: removed upload file {filename}")
                except Exception as e:
                    logger.error(f"Error removing upload file {filename}: {e}")
    
    # Check output folder
    if os.path.exists(OUTPUT_FOLDER):
        for filename in os.listdir(OUTPUT_FOLDER):
            file_path = os.path.join(OUTPUT_FOLDER, filename)
            # If file is older than 24 hours, delete it
            if os.path.isfile(file_path) and current_time - os.path.getmtime(file_path) > cleanup_time:
                try:
                    os.remove(file_path)
                    cleaned_files.append(f"output/{filename}")
                    logger.info(f"Auto cleanup: removed output file {filename}")
                except Exception as e:
                    logger.error(f"Error removing output file {filename}: {e}")
    
    # Clean up expired sessions (also after 24 hours)
    expired_sessions = []
    for session_id, session_data in sessions.items():
        if current_time - session_data.get('timestamp', 0) > cleanup_time:
            expired_sessions.append(session_id)
    
    for session_id in expired_sessions:
        try:
            del sessions[session_id]
            logger.info(f"Auto cleanup: removed expired session {session_id}")
        except Exception as e:
            logger.error(f"Error removing session {session_id}: {e}")
    
    if cleaned_files:
        logger.info(f"Auto cleanup completed: removed {len(cleaned_files)} files and {len(expired_sessions)} sessions")
    
    return len(cleaned_files), len(expired_sessions)

# Schedule cleanup to run every 24 hours
def run_cleanup_scheduler():
    # Run cleanup every 24 hours
    schedule.every(24).hours.do(cleanup_expired_files)
    
    # Also run cleanup at startup to clean any old files
    logger.info("Running initial cleanup at startup...")
    try:
        files_cleaned, sessions_cleaned = cleanup_expired_files()
        logger.info(f"Startup cleanup completed: {files_cleaned} files, {sessions_cleaned} sessions")
    except Exception as e:
        logger.error(f"Error during startup cleanup: {e}")
    
    while True:
        schedule.run_pending()
        time.sleep(3600)  # Check every hour, but cleanup runs every 24 hours

# Start cleanup thread
cleanup_thread = threading.Thread(target=run_cleanup_scheduler, daemon=True)
cleanup_thread.start()

# Routes
@app.route('/')
def index():
    return app.send_static_file('index.html')

@app.route('/api/upload/public', methods=['POST'])
@limiter.limit("15 per minute")
def upload_file_public():
    """Public upload endpoint without authentication for frontend use"""
    try:
        # Check if file is in request
        if 'file' not in request.files:
            logger.warning("Upload attempt without file")
            return jsonify({'error': 'No file part'}), 400
        
        file = request.files['file']
        
        # Check if file is selected
        if not file or file.filename == '':
            logger.warning("Upload attempt with empty filename")
            return jsonify({'error': 'No file selected'}), 400
        
        # Sanitize filename
        original_filename = file.filename
        sanitized_filename = sanitize_filename(original_filename)
        
        if not sanitized_filename:
            logger.warning(f"Invalid filename: {original_filename}")
            return jsonify({'error': 'Invalid filename'}), 400
        
        # Check if file type is allowed
        if not allowed_file(sanitized_filename):
            logger.warning(f"Unsupported file type: {sanitized_filename}")
            return jsonify({'error': 'File type not supported'}), 400
        
        # Check file size
        file.seek(0, os.SEEK_END)
        file_size = file.tell()
        file.seek(0)
        
        if file_size == 0:
            logger.warning("Empty file upload attempt")
            return jsonify({'error': 'File is empty'}), 400
        
        if file_size > MAX_FILE_SIZE:
            logger.warning(f"File too large: {file_size} bytes")
            return jsonify({'error': f'File too large. Maximum size is {MAX_FILE_SIZE // (1024*1024)}MB'}), 400
        
        # Get session ID from form data or create new session
        session_id = request.form.get('sessionId')
        if not session_id:
            session_id = str(uuid.uuid4())
            sessions[session_id] = {
                'created_at': time.time(),
                'files': {}
            }
            logger.info(f"Created new public session: {session_id}")
        elif session_id not in sessions:
            # Create session if it doesn't exist
            sessions[session_id] = {
                'created_at': time.time(),
                'files': {}
            }
            logger.info(f"Created missing public session: {session_id}")
        
        # Generate unique filename
        file_extension = get_file_extension(sanitized_filename)
        unique_filename = f"{uuid.uuid4()}.{file_extension}"
        file_path = os.path.join(UPLOAD_FOLDER, unique_filename)
        
        # Save file
        file.save(file_path)
        
        # Validate file content
        validation_result = validate_file_content(file_path, file_extension)
        if not validation_result['valid']:
            # Remove invalid file
            try:
                os.remove(file_path)
            except:
                pass
            logger.warning(f"Invalid file content: {validation_result['error']}")
            return jsonify({'error': validation_result['error']}), 400
        
        # Generate file ID
        file_id = str(uuid.uuid4())
        
        # Store file info in session
        session = sessions[session_id]
        session['files'][file_id] = {
            'path': file_path,
            'original_name': original_filename,
            'size': file_size,
            'upload_time': time.time(),
            'type': 'input'
        }
        # Explicitly save the session to trigger persistence
        sessions[session_id] = session
        
        logger.info(f"Public file uploaded successfully: {file_id} ({original_filename})")
        
        return jsonify({
            'success': True,
            'fileId': file_id,
            'sessionId': session_id,
            'filename': original_filename,
            'size': file_size
        })
        
    except Exception as e:
        logger.error(f"Unexpected error in public upload: {str(e)}")
        return jsonify({'error': 'Internal server error'}), 500

@app.route('/api/upload', methods=['POST'])
@limiter.limit("10 per minute")
def upload_file():
    # Verify user authentication
    user_data = verify_supabase_token(request)
    if not user_data:
        return jsonify({'error': 'Authentication required'}), 401
    """Upload file with comprehensive security validation"""
    try:
        # Check if file is in request
        if 'file' not in request.files:
            logger.warning("Upload attempt without file")
            return jsonify({'error': 'No file part'}), 400
        
        file = request.files['file']
        
        # Check if file is selected
        if not file or file.filename == '':
            logger.warning("Upload attempt with empty filename")
            return jsonify({'error': 'No file selected'}), 400
        
        # Sanitize filename
        original_filename = file.filename
        sanitized_filename = sanitize_filename(original_filename)
        
        if not sanitized_filename:
            logger.warning(f"Invalid filename: {original_filename}")
            return jsonify({'error': 'Invalid filename'}), 400
        
        # Check if file type is allowed
        if not allowed_file(sanitized_filename):
            logger.warning(f"Unsupported file type: {sanitized_filename}")
            return jsonify({'error': 'File type not supported'}), 400
        
        # Check file size
        file.seek(0, os.SEEK_END)
        file_size = file.tell()
        file.seek(0)
        
        if file_size == 0:
            logger.warning("Empty file upload attempt")
            return jsonify({'error': 'File is empty'}), 400
        
        if file_size > MAX_FILE_SIZE:
            logger.warning(f"File too large: {file_size} bytes")
            return jsonify({'error': f'File too large. Maximum size is {MAX_FILE_SIZE // (1024*1024)}MB'}), 400
        
        # Generate unique filename to prevent conflicts
        file_extension = get_file_extension(sanitized_filename)
        unique_filename = f"{str(uuid.uuid4())[:8]}_{sanitized_filename}"
        file_path = os.path.join(UPLOAD_FOLDER, unique_filename)
        
        # Save the file
        try:
            file.save(file_path)
        except Exception as e:
            logger.error(f"Failed to save file: {str(e)}")
            return jsonify({'error': 'Failed to save file'}), 500
        
        # Validate file content matches extension
        if not validate_file_content(file_path, file_extension):
            logger.warning(f"File content validation failed for: {sanitized_filename}")
            # Clean up the uploaded file
            try:
                os.remove(file_path)
            except:
                pass
            return jsonify({'error': 'File content does not match extension'}), 400
        
        # Generate a file ID
        file_id = str(uuid.uuid4())
        
        # Get session ID from request or generate a new one
        session_id = request.form.get('sessionId')
        if not session_id or not re.match(r'^[a-zA-Z0-9-_]+$', session_id):
            session_id = str(uuid.uuid4())
        
        # Store file info in session
        if session_id not in sessions:
            sessions[session_id] = {
                'files': {},
                'timestamp': time.time(),
                'upload_count': 0,
                'user_id': user_data.get('sub'),
                'user_email': user_data.get('email')
            }
        
        # Limit files per session - but clean up old files first
        session_files = sessions[session_id]['files']
        if len(session_files) >= 50:
            # Clean up oldest files when approaching limit
            current_time = time.time()
            files_to_remove = []
            
            # Sort files by timestamp and remove oldest ones
            sorted_files = sorted(session_files.items(), key=lambda x: x[1]['timestamp'])
            files_to_remove = sorted_files[:25]  # Remove oldest 25 files
            
            for file_id, file_info in files_to_remove:
                try:
                    if os.path.exists(file_info['path']):
                        os.remove(file_info['path'])
                    del session_files[file_id]
                    logger.info(f"Cleaned up old file from session: {file_id}")
                except Exception as e:
                    logger.error(f"Error cleaning up file {file_id}: {str(e)}")
            
            # If still too many files after cleanup, return error
            if len(session_files) >= 50:
                logger.warning(f"Too many files in session after cleanup: {session_id}")
                os.remove(file_path)  # Clean up current file
                return jsonify({'error': 'Session storage full. Please start a new session.'}), 400
        
        sessions[session_id]['files'][file_id] = {
            'original_name': sanitized_filename,
            'path': file_path,
            'size': file_size,
            'type': file_extension,
            'timestamp': time.time(),
            'upload_ip': get_remote_address()
        }
        
        # Update session timestamp and upload count
        sessions[session_id]['timestamp'] = time.time()
        sessions[session_id]['upload_count'] = sessions[session_id].get('upload_count', 0) + 1
        
        logger.info(f"File uploaded successfully: {sanitized_filename} ({file_size} bytes)")
        
        # Return file ID and session ID
        return jsonify({
            'fileId': file_id,
            'sessionId': session_id,
            'name': sanitized_filename,
            'size': file_size,
            'type': file_extension
        })
        
    except Exception as e:
        logger.error(f"Upload error: {str(e)}")
        return jsonify({'error': 'Internal server error'}), 500

@app.route('/api/convert/public', methods=['POST'])
@limiter.limit("10 per minute")
def convert_file_public():
    """Public conversion endpoint without authentication for frontend use"""
    try:
        # Handle malformed JSON
        try:
            data = request.json
        except Exception as e:
            logger.warning(f"Malformed JSON in convert request: {e}")
            return jsonify({'error': 'Invalid JSON format'}), 400
        
        # Validate request data
        if not data or 'fileId' not in data or 'outputFormat' not in data or 'sessionId' not in data:
            logger.warning("Convert request missing required parameters")
            return jsonify({'error': 'Missing required parameters: fileId, outputFormat, sessionId'}), 400
        
        file_id = data['fileId']
        output_format = data['outputFormat'].lower().strip()  # Normalize format to lowercase
        session_id = data['sessionId']
        options = data.get('options', {})
        
        # Validate input parameters
        if not re.match(r'^[a-zA-Z0-9-_]+$', file_id):
            logger.warning(f"Invalid file ID format: {file_id}")
            return jsonify({'error': 'Invalid file ID'}), 400
        
        if not re.match(r'^[a-zA-Z0-9-_]+$', session_id):
            logger.warning(f"Invalid session ID format: {session_id}")
            return jsonify({'error': 'Invalid session ID'}), 400
        
        # Use enhanced conversion system if available
        if conversion_manager:
            # Validate output format using enhanced system
            supported_formats = conversion_manager.get_supported_formats()
            if output_format not in supported_formats['outputs']:
                logger.warning(f"Unsupported output format: {output_format}")
                return jsonify({'error': f'Unsupported output format: {output_format}'}), 400
        else:
            # Fallback validation
            allowed_output_formats = ['pdf', 'docx', 'png', 'jpg', 'jpeg', 'txt', 'html', 'csv']
            if output_format not in allowed_output_formats:
                logger.warning(f"Unsupported output format: {output_format}")
                return jsonify({'error': f'Unsupported output format: {output_format}'}), 400
        
        # Check if session exists
        if session_id not in sessions:
            logger.warning(f"Session not found: {session_id}")
            return jsonify({'error': 'Session not found'}), 404
        
        session = sessions[session_id]
        
        # Check if file exists in session
        if file_id not in session['files']:
            logger.warning(f"File not found in session: {file_id}")
            return jsonify({'error': 'File not found'}), 404
        
        file_info = session['files'][file_id]
        input_path = file_info['path']
        
        # Check if input file exists
        if not os.path.exists(input_path):
            logger.warning(f"Input file not found: {input_path}")
            return jsonify({'error': 'Input file not found'}), 404
        
        # Get input file type
        input_type = get_file_extension(file_info['original_name']).lower()
        
        # Check file size
        if file_info['size'] > MAX_FILE_SIZE:
            logger.warning(f"File too large for conversion: {file_info['size']} bytes")
            return jsonify({'error': 'File too large for conversion'}), 400
        
        # Check if conversion is supported (enhanced system or fallback)
        conversion_supported = False
        use_enhanced_system = False
        
        if conversion_manager:
            # Check if enhanced system can handle this conversion
            if conversion_manager.can_convert(input_type, output_format):
                conversion_supported = True
                use_enhanced_system = True
                logger.info(f"Using enhanced conversion system for {input_type} to {output_format}")
            else:
                logger.info(f"Enhanced system cannot convert {input_type} to {output_format}, checking fallback")
        
        if not conversion_supported:
            # Check fallback conversion map
            conversion_map = {
                'pdf': ['docx', 'png', 'jpg', 'jpeg', 'txt', 'html'],
                'png': ['pdf', 'jpg', 'jpeg'],
                'jpg': ['pdf', 'png'],
                'jpeg': ['pdf', 'png'],
                'docx': ['pdf', 'txt'],
                'txt': ['pdf', 'docx', 'html']
            }
            
            if input_type in conversion_map and output_format in conversion_map[input_type]:
                conversion_supported = True
                use_enhanced_system = False
                logger.info(f"Using fallback conversion for {input_type} to {output_format}")
        
        if not conversion_supported:
            logger.warning(f"Unsupported conversion: {input_type} to {output_format}")
            return jsonify({'error': f'Cannot convert {input_type} to {output_format}'}), 400
        
        # Prevent converting to same format
        if input_type == output_format:
            logger.warning(f"Same format conversion attempted: {input_type}")
            return jsonify({'error': 'Source and target formats are the same'}), 400
        
        # Generate output filename
        output_filename = generate_unique_filename(file_info['original_name'], output_format)
        output_path = os.path.join(OUTPUT_FOLDER, output_filename)
        
        logger.info(f"Starting public conversion: {input_type} to {output_format}")
        
        # Use the appropriate conversion system
        conversion_successful = False
        if use_enhanced_system and conversion_manager:
            try:
                # Use the enhanced conversion manager
                result = conversion_manager.convert_file(
                    input_path=input_path,
                    output_path=output_path,
                    input_format=input_type,
                    output_format=output_format,
                    options=options
                )
                
                if result['success']:
                    logger.info(f"Enhanced conversion successful: {input_type} to {output_format}")
                    conversion_successful = True
                else:
                    logger.warning(f"Enhanced conversion failed: {result.get('error', 'Unknown error')}, trying fallback methods")
                
            except Exception as e:
                logger.warning(f"Enhanced conversion error: {str(e)}, trying fallback methods")
        
        # Use fallback conversion logic if enhanced system is not used or failed
        if not conversion_successful:
            # Fallback conversion logic
            try:
                if input_type == 'pdf' and output_format == 'docx':
                    convert_pdf_to_docx(input_path, output_path)
                elif input_type == 'pdf' and output_format in ['png', 'jpg', 'jpeg']:
                    convert_pdf_to_image(input_path, output_path, output_format)
                elif input_type == 'pdf' and output_format == 'txt':
                    convert_pdf_to_text(input_path, output_path)
                elif input_type in ['png', 'jpg', 'jpeg'] and output_format == 'pdf':
                    convert_image_to_pdf(input_path, output_path)
                elif input_type in ['png', 'jpg', 'jpeg'] and output_format in ['png', 'jpg', 'jpeg']:
                    convert_image_to_image(input_path, output_path, output_format)
                elif input_type == 'docx' and output_format == 'pdf':
                    convert_docx_to_pdf(input_path, output_path)
                elif input_type == 'docx' and output_format == 'txt':
                    convert_docx_to_text(input_path, output_path)
                elif input_type == 'txt' and output_format == 'pdf':
                    convert_text_to_pdf(input_path, output_path)
                elif input_type == 'txt' and output_format == 'html':
                    convert_text_to_html(input_path, output_path)
                else:
                    logger.warning(f"Unsupported conversion: {input_type} to {output_format}")
                    return jsonify({'error': f'Conversion from {input_type} to {output_format} not supported'}), 400
                
                logger.info(f"Fallback conversion successful: {input_type} to {output_format}")
                
            except Exception as e:
                logger.error(f"Conversion error: {str(e)}")
                return jsonify({'error': f'Conversion failed: {str(e)}'}), 500
        
        # Check if output file was created
        if not os.path.exists(output_path):
            logger.error(f"Output file not created: {output_path}")
            return jsonify({'error': 'Conversion failed - output file not created'}), 500
        
        # Get output file size
        output_size = os.path.getsize(output_path)
        
        # Generate unique ID for output file
        output_file_id = str(uuid.uuid4())
        
        # Store output file info in session
        session['files'][output_file_id] = {
            'path': output_path,
            'original_name': output_filename,
            'size': output_size,
            'upload_time': time.time(),
            'type': 'output'
        }
        # Explicitly save the session to trigger persistence
        sessions[session_id] = session
        
        logger.info(f"Public conversion completed successfully: {file_id} -> {output_file_id}")
        
        return jsonify({
            'success': True,
            'fileId': output_file_id,
            'filename': output_filename,
            'size': output_size,
            'downloadUrl': f'/api/download/public/{output_file_id}?session_id={session_id}'
        })
        
    except Exception as e:
        logger.error(f"Unexpected error in public conversion: {str(e)}")
        return jsonify({'error': 'Internal server error'}), 500

@app.route('/api/convert', methods=['POST'])
@limiter.limit("5 per minute")
def convert_file():
    # Verify user authentication
    user_data = verify_supabase_token(request)
    if not user_data:
        return jsonify({'error': 'Authentication required'}), 401
    """Convert file with comprehensive validation and security checks"""
    try:
        # Handle malformed JSON
        try:
            data = request.json
        except Exception as e:
            logger.warning(f"Malformed JSON in convert request: {e}")
            return jsonify({'error': 'Invalid JSON format'}), 400
        
        # Validate request data
        if not data or 'fileId' not in data or 'outputFormat' not in data or 'sessionId' not in data:
            logger.warning("Convert request missing required parameters")
            return jsonify({'error': 'Missing required parameters: fileId, outputFormat, sessionId'}), 400
        
        file_id = data['fileId']
        output_format = data['outputFormat'].lower().strip()  # Normalize format to lowercase
        session_id = data['sessionId']
        options = data.get('options', {})
        
        # Validate input parameters
        if not re.match(r'^[a-zA-Z0-9-_]+$', file_id):
            logger.warning(f"Invalid file ID format: {file_id}")
            return jsonify({'error': 'Invalid file ID'}), 400
        
        if not re.match(r'^[a-zA-Z0-9-_]+$', session_id):
            logger.warning(f"Invalid session ID format: {session_id}")
            return jsonify({'error': 'Invalid session ID'}), 400
        
        # Use enhanced conversion system if available
        if conversion_manager:
            # Validate output format using enhanced system
            supported_formats = conversion_manager.get_supported_formats()
            if output_format not in supported_formats['outputs']:
                logger.warning(f"Unsupported output format: {output_format}")
                return jsonify({'error': 'Unsupported output format'}), 400
        else:
            # Fallback to basic format validation
            allowed_formats = ['pdf', 'docx', 'png', 'jpg', 'jpeg', 'txt']
            if output_format not in allowed_formats:
                logger.warning(f"Unsupported output format: {output_format}")
                return jsonify({'error': 'Unsupported output format'}), 400
        
        # Check if session exists
        if session_id not in sessions:
            logger.warning(f"Session not found: {session_id}")
            return jsonify({'error': 'Session not found'}), 404
        
        # Check if file exists in session
        if file_id not in sessions[session_id]['files']:
            logger.warning(f"File not found in session: {file_id}")
            return jsonify({'error': 'File not found'}), 404
        
        # Get file info
        file_info = sessions[session_id]['files'][file_id]
        input_path = file_info['path']
        input_type = file_info['type'].lower()  # Normalize type to lowercase
        
        # Check if file exists on disk
        if not os.path.exists(input_path):
            logger.error(f"File not found on disk: {input_path}")
            return jsonify({'error': 'File not found on disk'}), 404
        
        # Validate file size before conversion
        if file_info['size'] > MAX_FILE_SIZE:
            logger.warning(f"File too large for conversion: {file_info['size']} bytes")
            return jsonify({'error': 'File too large for conversion'}), 400
        
        # Use enhanced conversion system if available
        if conversion_manager:
            # Validate conversion using enhanced system
            if not conversion_manager.can_convert(input_type, output_format):
                logger.warning(f"Unsupported conversion: {input_type} to {output_format}")
                return jsonify({'error': f'Cannot convert {input_type} to {output_format}'}), 400
        else:
            # Fallback to basic conversion map
            conversion_map = {
                'pdf': ['docx', 'png', 'jpg', 'jpeg', 'txt', 'html'],
                'png': ['pdf', 'jpg', 'jpeg'],
                'jpg': ['pdf', 'png'],
                'jpeg': ['pdf', 'png'],
                'docx': ['pdf', 'txt'],
                'txt': ['pdf', 'docx']
            }
            
            if input_type not in conversion_map:
                logger.warning(f"Unsupported source format: {input_type}")
                return jsonify({'error': 'Unsupported source format'}), 400
            
            if output_format not in conversion_map[input_type]:
                logger.warning(f"Unsupported conversion: {input_type} to {output_format}")
                return jsonify({'error': f'Cannot convert {input_type} to {output_format}'}), 400
        
        # Prevent converting to same format
        if input_type == output_format:
            logger.warning(f"Same format conversion attempted: {input_type}")
            return jsonify({'error': 'Source and target formats are the same'}), 400
        
        # Generate output filename
        output_filename = generate_unique_filename(file_info['original_name'], output_format)
        output_path = os.path.join(OUTPUT_FOLDER, output_filename)
        
        logger.info(f"Starting conversion: {input_type} to {output_format}")
        
        # Use enhanced conversion system if available
        conversion_successful = False
        if conversion_manager:
            try:
                # Use the enhanced conversion manager
                result = conversion_manager.convert_file(
                    input_path=input_path,
                    output_path=output_path,
                    input_format=input_type,
                    output_format=output_format,
                    options=options
                )
                
                if result['success']:
                    # Update output path if the conversion manager changed it
                    if result.get('output_path'):
                        output_path = result['output_path']
                    conversion_successful = True
                    logger.info(f"Enhanced conversion successful: {input_type} to {output_format}")
                else:
                    logger.warning(f"Enhanced conversion failed: {result.get('error', 'Unknown error')}, trying fallback methods")
                    
            except Exception as conversion_error:
                logger.warning(f"Enhanced conversion error: {str(conversion_error)}, trying fallback methods")
        
        # Use fallback conversion logic if enhanced system is not available or failed
        if not conversion_successful:
            logger.info("Using legacy conversion methods")
            
            # Perform conversion based on input and output formats (legacy fallback)
            if input_type == 'pdf':
                if output_format == 'docx':
                    convert_pdf_to_docx(input_path, output_path)
                elif output_format in ['jpg', 'jpeg', 'png']:
                    convert_pdf_to_image(input_path, output_path, format=output_format)
                elif output_format == 'txt':
                    convert_pdf_to_text(input_path, output_path, use_ocr=options.get('ocr', False))
            elif input_type in ['jpg', 'jpeg', 'png']:
                if output_format == 'pdf':
                    convert_image_to_pdf(input_path, output_path)
                elif output_format in ['jpg', 'jpeg', 'png']:
                    convert_image_to_image(input_path, output_path, output_format)
            elif input_type == 'docx':
                if output_format == 'pdf':
                    convert_docx_to_pdf(input_path, output_path)
                elif output_format == 'txt':
                    convert_docx_to_text(input_path, output_path)
            elif input_type == 'txt':
                if output_format == 'pdf':
                    convert_text_to_pdf(input_path, output_path)
            else:
                raise Exception(f"Unsupported conversion: {input_type} to {output_format}")
        
        # Verify the output file exists and has content
        if not os.path.exists(output_path) or os.path.getsize(output_path) == 0:
            logger.error(f"Conversion produced empty or missing file: {output_path}")
            return jsonify({'error': 'Conversion failed: Output file is empty or not created'}), 500
        
        # Determine the correct MIME type for the output file
        mime_type = None
        if output_format == 'pdf':
            mime_type = 'application/pdf'
        elif output_format == 'docx':
            mime_type = 'application/vnd.openxmlformats-officedocument.wordprocessingml.document'
        elif output_format == 'xlsx':
            mime_type = 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet'
        elif output_format == 'pptx':
            mime_type = 'application/vnd.openxmlformats-officedocument.presentationml.presentation'
        elif output_format == 'csv':
            mime_type = 'text/csv'
        elif output_format == 'html':
            mime_type = 'text/html'
        elif output_format in ['jpg', 'jpeg']:
            mime_type = 'image/jpeg'
        elif output_format == 'png':
            mime_type = 'image/png'
        elif output_format == 'webp':
            mime_type = 'image/webp'
        elif output_format == 'bmp':
            mime_type = 'image/bmp'
        elif output_format == 'tiff':
            mime_type = 'image/tiff'
        elif output_format == 'gif':
            mime_type = 'image/gif'
        elif output_format == 'txt':
            mime_type = 'text/plain'
        
        # Store converted file info in session
        converted_file_id = str(uuid.uuid4())
        sessions[session_id]['files'][converted_file_id] = {
            'original_name': output_filename,
            'path': output_path,
            'size': os.path.getsize(output_path),
            'type': output_format,
            'mime_type': mime_type,  # Store the MIME type for later use
            'timestamp': time.time(),
            'converted_from': file_id,
            'conversion_ip': get_remote_address()
        }
        
        # Update session timestamp
        sessions[session_id]['timestamp'] = time.time()
        
        # Generate download URL
        download_url = url_for('download_file', file_id=converted_file_id, session_id=session_id, _external=True)
        
        logger.info(f"Conversion successful: {input_type} to {output_format} ({os.path.getsize(output_path)} bytes)")
        
        # Return success response
        return jsonify({
            'fileId': converted_file_id,
            'name': output_filename,
            'size': os.path.getsize(output_path),
            'type': output_format,
            'mimeType': mime_type,
            'url': download_url,
            'sessionUrl': url_for('index', session=session_id, _external=True)
        })
    
    except Exception as e:
        # Log the error
        logger.error(f"Conversion error: {str(e)}")
        # Return error response
        return jsonify({'error': 'Internal server error'}), 500

@app.route('/api/download/public/<file_id>', methods=['GET'])
@limiter.limit("30 per minute")
def download_file_public(file_id):
    """Public download endpoint for files created through public conversion flow"""
    session_id = request.args.get('session_id')
    # Check if this is a preview request
    is_preview = request.args.get('preview', 'false').lower() == 'true'
    # Check if this is a forced download request
    force_download_param = request.args.get('download', 'false').lower() == 'true'
    
    # Check if session exists
    if not session_id or session_id not in sessions:
        return jsonify({'error': 'Session not found'}), 404
    
    # Check if file exists in session
    if file_id not in sessions[session_id]['files']:
        return jsonify({'error': 'File not found'}), 404
    
    # Get file info
    file_info = sessions[session_id]['files'][file_id]
    file_path = file_info['path']
    file_type = file_info['type'].lower()
    
    # Check if file exists on disk
    if not os.path.exists(file_path):
        return jsonify({'error': 'File not found on server'}), 404
    
    # Update session timestamp
    sessions[session_id]['timestamp'] = time.time()
    
    # Use stored MIME type if available, otherwise determine based on file extension
    mime_type = file_info.get('mime_type')
    
    if mime_type is None:
        # Set the correct MIME type based on file extension
        if file_type == 'pdf':
            mime_type = 'application/pdf'
        elif file_type == 'docx':
            mime_type = 'application/vnd.openxmlformats-officedocument.wordprocessingml.document'
        elif file_type == 'xlsx':
            mime_type = 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet'
        elif file_type == 'pptx':
            mime_type = 'application/vnd.openxmlformats-officedocument.presentationml.presentation'
        elif file_type in ['jpg', 'jpeg']:
            mime_type = 'image/jpeg'
        elif file_type == 'png':
            mime_type = 'image/png'
        elif file_type == 'txt':
            mime_type = 'text/plain'
        elif file_type == 'html':
            mime_type = 'text/html'
        elif file_type == 'json':
            mime_type = 'application/json'
        
        # If no mime type was determined, try to guess based on file extension
        if mime_type is None:
            mime_type = 'application/octet-stream'  # Default fallback
        
    # For image files, verify they exist and are valid
    if file_type in ['jpg', 'jpeg', 'png'] and os.path.exists(file_path) and PIL_AVAILABLE:
        try:
            # Attempt to open the image to verify it's valid
            with Image.open(file_path) as img:
                # Get actual format from the image
                actual_format = img.format.lower() if img.format else file_type
                if actual_format == 'jpeg':
                    mime_type = 'image/jpeg'
                elif actual_format == 'png':
                    mime_type = 'image/png'
        except Exception as e:
            print(f"Error verifying image: {str(e)}")
            # Continue anyway, we'll try to serve the file
    
    # Determine if we should force download or display inline
    # Always display inline for preview requests
    # Force download if explicitly requested via download parameter
    # Otherwise, display images and PDFs inline, download other formats
    if is_preview:
        force_download = False  # Always inline for previews
    elif force_download_param:
        force_download = True   # Force download when explicitly requested
    else:
        force_download = not (file_type in ['jpg', 'jpeg', 'png', 'pdf'])  # Default behavior
    
    # Add cache control headers to improve performance
    response = send_file(file_path, 
                     as_attachment=force_download, 
                     download_name=file_info['original_name'], 
                     mimetype=mime_type)
    
    # Add cache headers for better performance
    response.headers['Cache-Control'] = 'public, max-age=300'  # Cache for 5 minutes
    logger.info(f"Public file downloaded: {file_id} ({file_info['original_name']})")
    return response

@app.route('/api/download/<file_id>', methods=['GET'])
def download_file(file_id):
    """Download endpoint with flexible authentication - supports both authenticated users and public sessions"""
    session_id = request.args.get('session_id')
    
    # Try to verify user authentication, but don't require it if session_id is provided
    user_data = verify_supabase_token(request)
    
    # If no user authentication and no session_id, require authentication
    if not user_data and not session_id:
        return jsonify({'error': 'Authentication required'}), 401
    
    # Check if this is a preview request
    is_preview = request.args.get('preview', 'false').lower() == 'true'
    # Check if this is a forced download request
    force_download_param = request.args.get('download', 'false').lower() == 'true'
    
    # Check if session exists
    if not session_id or session_id not in sessions:
        return jsonify({'error': 'Session not found'}), 404
    
    # Check if file exists in session
    if file_id not in sessions[session_id]['files']:
        return jsonify({'error': 'File not found'}), 404
    
    # Get file info
    file_info = sessions[session_id]['files'][file_id]
    file_path = file_info['path']
    file_type = file_info['type'].lower()
    
    # Check if file exists on disk
    if not os.path.exists(file_path):
        return jsonify({'error': 'File not found on server'}), 404
    
    # Update session timestamp
    sessions[session_id]['timestamp'] = time.time()
    
    # Use stored MIME type if available, otherwise determine based on file extension
    mime_type = file_info.get('mime_type')
    
    if mime_type is None:
        # Set the correct MIME type based on file extension
        if file_type == 'pdf':
            mime_type = 'application/pdf'
        elif file_type == 'docx':
            mime_type = 'application/vnd.openxmlformats-officedocument.wordprocessingml.document'
        elif file_type == 'xlsx':
            mime_type = 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet'
        elif file_type == 'pptx':
            mime_type = 'application/vnd.openxmlformats-officedocument.presentationml.presentation'
        elif file_type in ['jpg', 'jpeg']:
            mime_type = 'image/jpeg'
        elif file_type == 'png':
            mime_type = 'image/png'
        elif file_type == 'txt':
            mime_type = 'text/plain'
        elif file_type == 'html':
            mime_type = 'text/html'
        elif file_type == 'json':
            mime_type = 'application/json'
        
        # If no mime type was determined, try to guess based on file extension
        if mime_type is None:
            mime_type = 'application/octet-stream'  # Default fallback
        
    # For image files, verify they exist and are valid
    if file_type in ['jpg', 'jpeg', 'png'] and os.path.exists(file_path) and PIL_AVAILABLE:
        try:
            # Attempt to open the image to verify it's valid
            with Image.open(file_path) as img:
                # Get actual format from the image
                actual_format = img.format.lower() if img.format else file_type
                if actual_format == 'jpeg':
                    mime_type = 'image/jpeg'
                elif actual_format == 'png':
                    mime_type = 'image/png'
        except Exception as e:
            print(f"Error verifying image: {str(e)}")
            # Continue anyway, we'll try to serve the file
    
    # Determine if we should force download or display inline
    # Always display inline for preview requests
    # Force download if explicitly requested via download parameter
    # Otherwise, display images and PDFs inline, download other formats
    if is_preview:
        force_download = False  # Always inline for previews
    elif force_download_param:
        force_download = True   # Force download when explicitly requested
    else:
        force_download = not (file_type in ['jpg', 'jpeg', 'png', 'pdf'])  # Default behavior
    
    # Add cache control headers to improve performance
    response = send_file(file_path, 
                     as_attachment=force_download, 
                     download_name=file_info['original_name'], 
                     mimetype=mime_type)
    
    # Add cache headers for better performance
    response.headers['Cache-Control'] = 'public, max-age=300'  # Cache for 5 minutes
    return response

@app.route('/api/session/<session_id>', methods=['GET'])
def get_session(session_id):
    # Check if session exists
    if session_id not in sessions:
        return jsonify({'error': 'Session not found'}), 404
    
    # Get session info
    session_info = sessions[session_id]
    
    # Filter out sensitive info and prepare response
    files = []
    for file_id, file_info in session_info['files'].items():
        files.append({
            'id': file_id,
            'name': file_info['original_name'],
            'size': file_info['size'],
            'type': file_info['type'],
            'timestamp': file_info['timestamp'],
            'url': url_for('download_file', file_id=file_id, session_id=session_id, _external=True) if os.path.exists(file_info['path']) else None
        })
    
    # Sort files by timestamp (newest first)
    files.sort(key=lambda x: x['timestamp'], reverse=True)
    
    return jsonify({
        'sessionId': session_id,
        'files': files,
        'expiresAt': session_info['timestamp'] + FILE_EXPIRY
    })

@app.route('/api/session/<session_id>/reset', methods=['POST'])
def reset_session(session_id):
    """Reset a session by clearing all its files"""
    if session_id not in sessions:
        return jsonify({'error': 'Session not found'}), 404
    
    session_info = sessions[session_id]
    files_removed = 0
    
    # Remove all files from the session
    for file_id, file_info in session_info['files'].items():
        try:
            if os.path.exists(file_info['path']):
                os.remove(file_info['path'])
                files_removed += 1
        except Exception as e:
            logger.error(f"Error removing file {file_id}: {str(e)}")
    
    # Clear the files dictionary but keep the session
    sessions[session_id]['files'] = {}
    sessions[session_id]['timestamp'] = time.time()  # Update timestamp
    sessions[session_id]['upload_count'] = 0  # Reset upload count
    
    logger.info(f"Reset session {session_id}, removed {files_removed} files")
    
    return jsonify({
        'message': 'Session reset successfully',
        'filesRemoved': files_removed,
        'sessionId': session_id
    })

# Error handlers
@app.errorhandler(404)
def not_found(error):
    return app.send_static_file('index.html')

@app.errorhandler(413)
def request_entity_too_large(error):
    return jsonify({'error': 'File too large'}), 413

@app.errorhandler(500)
def server_error(error):
    return jsonify({'error': 'Server error'}), 500

# Malformed JSON error handler
@app.errorhandler(400)
def bad_request(error):
    return jsonify({'error': 'Bad request', 'message': 'Invalid request format or missing required fields'}), 400

# Route to inject environment variables into frontend
@app.route('/api/config')
def get_config():
    """Provide frontend configuration without exposing sensitive data"""
    return jsonify({
        'supabaseUrl': SUPABASE_URL,
        'supabaseAnonKey': SUPABASE_ANON_KEY,
        'maxFileSize': MAX_FILE_SIZE,
        'allowedExtensions': list(ALLOWED_EXTENSIONS)
    })

@app.route('/api/conversion/capabilities', methods=['GET'])
def get_conversion_capabilities():
    """Get enhanced conversion capabilities and format matrix"""
    try:
        if not conversion_manager:
            return jsonify({
                'error': 'Enhanced conversion system not available',
                'fallback': True
            }), 503
        
        # Get comprehensive conversion matrix
        conversion_matrix = conversion_manager.get_conversion_matrix()
        supported_formats = conversion_manager.get_supported_formats()
        engine_status = conversion_manager.get_engine_status()
        
        return jsonify({
            'conversionMatrix': conversion_matrix,
            'supportedFormats': supported_formats,
            'engines': engine_status,
            'enhanced': True,
            'version': '2.0.0'
        })
        
    except Exception as e:
        logger.error(f"Error getting conversion capabilities: {str(e)}")
        return jsonify({'error': 'Failed to get conversion capabilities'}), 500

@app.route('/api/conversion/options/<input_format>/<output_format>', methods=['GET'])
def get_conversion_options(input_format, output_format):
    """Get available conversion options for specific format pair"""
    try:
        if not conversion_manager:
            return jsonify({'error': 'Enhanced conversion system not available'}), 503
        
        options = conversion_manager.get_conversion_options(input_format, output_format)
        format_info = {
            'input': conversion_manager.get_format_info(input_format),
            'output': conversion_manager.get_format_info(output_format)
        }
        
        return jsonify({
            'options': options,
            'formatInfo': format_info
        })
        
    except Exception as e:
        logger.error(f"Error getting conversion options: {str(e)}")
        return jsonify({'error': 'Failed to get conversion options'}), 500

@app.route('/health', methods=['GET'])
def health_check():
    """Simple health check endpoint for production monitoring"""
    try:
        # Basic health checks
        health_status = {
            'status': 'healthy',
            'timestamp': datetime.now(timezone.utc).isoformat(),
            'version': '1.0.0',
            'uptime': time.time() - app.start_time if hasattr(app, 'start_time') else 'unknown'
        }
        
        # Check if upload and output directories exist
        if not os.path.exists(UPLOAD_FOLDER):
            health_status['status'] = 'degraded'
            health_status['issues'] = health_status.get('issues', []) + ['Upload folder missing']
        
        if not os.path.exists(OUTPUT_FOLDER):
            health_status['status'] = 'degraded'
            health_status['issues'] = health_status.get('issues', []) + ['Output folder missing']
        
        # Check conversion system availability
        health_status['conversion_system'] = conversion_system_available
        
        return jsonify(health_status), 200 if health_status['status'] == 'healthy' else 503
        
    except Exception as e:
        logger.error(f"Health check failed: {str(e)}")
        return jsonify({
            'status': 'unhealthy',
            'error': 'Health check failed',
            'timestamp': datetime.now(timezone.utc).isoformat()
        }), 503

# =============================================================================
# CONTACT FORM HANDLER
# =============================================================================

@app.route('/api/contact', methods=['POST'])
@limiter.limit("5 per minute")
def contact_form():
    """Handle contact form submissions and route to contact@entityconsulting.net"""
    try:
        data = request.get_json()
        
        # Validate required fields
        required_fields = ['name', 'email', 'subject', 'message']
        for field in required_fields:
            if not data.get(field):
                return jsonify({
                    'success': False,
                    'error': f'Missing required field: {field}'
                }), 400
        
        # Validate email format
        if not EMAIL_PATTERN.match(data['email']):
            return jsonify({
                'success': False,
                'error': 'Invalid email format'
            }), 400
        
        # Sanitize inputs
        name = data['name'].strip()[:100]
        email = data['email'].strip()[:100]
        subject = data['subject'].strip()[:200]
        message = data['message'].strip()[:2000]
        inquiry_type = data.get('inquiry_type', 'general').strip()[:50]
        
        # Create email content
        email_subject = f"[DocSwap Contact] {inquiry_type.title()}: {subject}"
        email_body = f"""
New contact form submission from DocSwap:

Name: {name}
Email: {email}
Inquiry Type: {inquiry_type.title()}
Subject: {subject}

Message:
{message}

---
Submitted from: {get_remote_address()}
Timestamp: {datetime.now(timezone.utc).isoformat()}
"""
        
        # For now, log the contact form submission
        # In production, you would integrate with your email service
        logger.info(f"Contact form submission: {email_subject}")
        logger.info(f"From: {name} <{email}>")
        logger.info(f"Message: {message[:200]}...")
        
        # Store contact submission for admin review
        contact_data = {
            'id': str(uuid.uuid4()),
            'name': name,
            'email': email,
            'subject': subject,
            'message': message,
            'inquiry_type': inquiry_type,
            'timestamp': datetime.now(timezone.utc).isoformat(),
            'ip_address': get_remote_address(),
            'status': 'new'
        }
        
        # Save to sessions for admin review (you can implement database storage later)
        contact_id = f"contact_{contact_data['id']}"
        sessions[contact_id] = contact_data
        
        return jsonify({
            'success': True,
            'message': 'Thank you for your message! We will get back to you within 24 hours.',
            'contact_id': contact_data['id']
        }), 200
        
    except Exception as e:
        logger.error(f"Contact form error: {str(e)}")
        return jsonify({
            'success': False,
            'error': 'Failed to send message. Please try again later.'
        }), 500

@app.route('/api/contact/admin', methods=['GET'])
def get_contact_submissions():
    """Get all contact form submissions for admin review"""
    try:
        # This would require admin authentication in production
        contact_submissions = []
        for key in sessions.keys():
            if key.startswith('contact_'):
                contact_submissions.append(sessions[key])
        
        # Sort by timestamp (newest first)
        contact_submissions.sort(key=lambda x: x['timestamp'], reverse=True)
        
        return jsonify({
            'success': True,
            'submissions': contact_submissions
        }), 200
        
    except Exception as e:
        logger.error(f"Error fetching contact submissions: {str(e)}")
        return jsonify({
            'success': False,
            'error': 'Failed to fetch contact submissions'
        }), 500

# Import admin module
from admin import admin_app, init_admin

# Register admin blueprint
app.register_blueprint(admin_app, url_prefix='/admin')

# Initialize admin interface
init_admin(app)

if __name__ == '__main__':
    # Disable debug mode in production
    debug_mode = os.getenv('FLASK_DEBUG', 'False').lower() == 'true'
    host = os.getenv('HOST', '0.0.0.0')  # Use 0.0.0.0 for production hosting
    port = int(os.getenv('PORT', 8000))
    app.run(debug=debug_mode, host=host, port=port)